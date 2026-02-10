#!/usr/bin/env python3
"""
Generate the Pluto Jacks 30% Capstone Review PowerPoint Presentation.

NAU ASCE Concrete Canoe 2026
Output: /root/concrete-canoe-project2026/presentation/Pluto_Jacks_30_Percent_Review.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# =============================================================================
# CONSTANTS
# =============================================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(0x1B, 0x36, 0x5D)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x88, 0x00)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)

FIGURES_DIR = "/root/concrete-canoe-project2026/reports/figures"
OUTPUT_DIR = "/root/concrete-canoe-project2026/presentation"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Pluto_Jacks_30_Percent_Review.pptx")


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================
def set_slide_bg(slide, color):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", italic=False):
    """Add a simple text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=18, color=BLACK, bold=False,
                          alignment=PP_ALIGN.LEFT, font_name="Calibri",
                          line_spacing=1.2):
    """Add a text box with multiple lines (each line as a separate paragraph)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, dict):
            p.text = line_data.get("text", "")
            p.font.size = Pt(line_data.get("size", font_size))
            p.font.color.rgb = line_data.get("color", color)
            p.font.bold = line_data.get("bold", bold)
            p.font.italic = line_data.get("italic", False)
            p.font.name = line_data.get("font_name", font_name)
            p.alignment = line_data.get("alignment", alignment)
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = font_name
            p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1.0) + 2)
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=18, color=BLACK, font_name="Calibri",
                    bullet_char="\u2022", line_spacing=1.3):
    """Add a bulleted list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, dict):
            p.text = f"{bullet_char} {item.get('text', '')}"
            p.font.size = Pt(item.get("size", font_size))
            p.font.color.rgb = item.get("color", color)
            p.font.bold = item.get("bold", False)
            p.font.name = font_name
        else:
            p.text = f"{bullet_char} {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1.0) + 2)
    return txBox


def add_slide_number(slide, number, total=24, color=MEDIUM_GRAY):
    """Add slide number in bottom-right corner."""
    add_textbox(
        slide,
        left=Inches(12.0), top=Inches(7.0),
        width=Inches(1.2), height=Inches(0.4),
        text=f"{number} / {total}",
        font_size=12, color=color,
        alignment=PP_ALIGN.RIGHT
    )


def add_speaker_tag(slide, speaker, color=MEDIUM_GRAY):
    """Add 'Speaker: Name' in the bottom-left corner."""
    add_textbox(
        slide,
        left=Inches(0.3), top=Inches(7.0),
        width=Inches(3.0), height=Inches(0.4),
        text=f"Speaker: {speaker}",
        font_size=12, color=color, italic=True
    )


def add_notes(slide, text):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_title_bar(slide, title_text, subtitle_text=None, bg_color=NAVY):
    """Add a navy blue title bar at the top of a content slide."""
    # Title bar background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    # Title text
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.15),
        width=Inches(12.0), height=Inches(0.7),
        text=title_text,
        font_size=32, color=WHITE, bold=True,
        font_name="Calibri"
    )

    if subtitle_text:
        add_textbox(
            slide,
            left=Inches(0.5), top=Inches(0.75),
            width=Inches(12.0), height=Inches(0.4),
            text=subtitle_text,
            font_size=16, color=GOLD, bold=False,
            font_name="Calibri", italic=True
        )


def add_gold_accent_line(slide, top):
    """Add a thin gold accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), top, Inches(12.333), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add an image to the slide, or a placeholder if file not found."""
    if os.path.exists(image_path):
        kwargs = {"image_file": image_path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        try:
            slide.shapes.add_picture(**kwargs)
            return True
        except Exception as e:
            print(f"  Warning: Could not load image {image_path}: {e}")
    else:
        print(f"  Warning: Image not found: {image_path}")

    # Placeholder
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top,
        width or Inches(4), height or Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = MEDIUM_GRAY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[Image: {os.path.basename(image_path)}]"
    p.font.size = Pt(14)
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.CENTER
    return False


def add_table(slide, left, top, width, height, rows, cols, data,
              header_bg=NAVY, header_fg=WHITE, body_fg=BLACK,
              col_widths=None, font_size=14):
    """
    Add a formatted table.
    data: list of lists. First row is header.
    """
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER

                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_fg
                else:
                    paragraph.font.color.rgb = body_fg

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEB, 0xF0)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_footer_text(slide, text, color=MEDIUM_GRAY):
    """Add a footer text line above the slide number area."""
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(6.7),
        width=Inches(12.0), height=Inches(0.3),
        text=text,
        font_size=10, color=color, italic=True,
        alignment=PP_ALIGN.LEFT
    )


# =============================================================================
# SLIDE BUILDERS
# =============================================================================
def build_slide_01_title(prs):
    """Slide 1: Title Slide"""
    print("  Building Slide 1: Title")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # Gold accent bar near top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(1.0), Inches(9.333), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    # PLUTO JACKS
    add_textbox(
        slide, Inches(0.5), Inches(1.3), Inches(12.333), Inches(1.2),
        "PLUTO JACKS", font_size=54, color=WHITE, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Subtitle
    add_textbox(
        slide, Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.7),
        "NAU ASCE Concrete Canoe 2026", font_size=28, color=GOLD,
        bold=False, alignment=PP_ALIGN.CENTER
    )

    # Capstone line
    add_textbox(
        slide, Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.6),
        "CENE 476 Capstone  --  30% Design Review", font_size=22, color=WHITE,
        alignment=PP_ALIGN.CENTER
    )

    # Gold accent bar
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.0), Inches(3.9), Inches(7.333), Inches(0.04)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GOLD
    shape2.line.fill.background()

    # Date
    add_textbox(
        slide, Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5),
        "February 10, 2026", font_size=20, color=WHITE,
        alignment=PP_ALIGN.CENTER
    )

    # Team
    add_textbox(
        slide, Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.5),
        "Team:  Trevion  |  Kayleigh  |  Amit  |  Jon  |  Theo",
        font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER
    )

    # Advisors
    add_textbox(
        slide, Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5),
        "Advisors:  Dr. Lamer PE  |  Taylor Layland PE",
        font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER, italic=True
    )

    add_speaker_tag(slide, "Theo", color=RGBColor(0x99, 0x99, 0x99))
    add_slide_number(slide, 1, color=RGBColor(0x99, 0x99, 0x99))
    add_notes(slide, "We're Pluto Jacks, NAU's 2025-2026 ASCE Concrete Canoe team. Five members. Today we present our 30% progress.")


def build_slide_02_introduction(prs):
    """Slide 2: Introduction & Competition"""
    print("  Building Slide 2: Introduction")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Project Introduction")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Design and build a concrete canoe per ASCE 2026 rules",
        "National competition: Salt Lake City, University of Utah",
        "Four event categories: Technical Paper, Oral Presentation, Final Product, Racing",
        "Team of 5 undergraduate civil engineers at Northern Arizona University",
        "Faculty Advisor: Dr. Lamer PE",
        "Graduate Instructor: Taylor Layland PE",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(7.5), Inches(4.5),
        items, font_size=20, color=DARK_GRAY
    )

    # Team photo placeholder box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.8), Inches(1.8), Inches(4.0), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = MEDIUM_GRAY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Team Photo Placeholder]"
    p.font.size = Pt(16)
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 2)
    add_notes(slide, "Purpose: design and build a concrete canoe per ASCE rules. Competing in Salt Lake City at University of Utah. Five-member team from NAU Civil Engineering.")


def build_slide_03_schedule(prs):
    """Slide 3: Schedule & Milestones"""
    print("  Building Slide 3: Schedule")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Project Milestones")
    add_gold_accent_line(slide, Inches(1.25))

    # Status summary
    data = [
        ["Task", "Description", "Status", "Progress"],
        ["Task 1", "Project Research & Safety", "COMPLETE", "100%"],
        ["Task 2", "Material Testing & Mix Selection", "COMPLETE", "100%"],
        ["Task 3", "Hull Design & Structural Analysis", "IN PROGRESS", "60%"],
        ["Task 4", "Construction & Fabrication", "NOT STARTED", "0%"],
        ["Task 5", "Finishing & Competition Prep", "NOT STARTED", "0%"],
    ]
    add_table(
        slide,
        left=Inches(0.5), top=Inches(1.5),
        width=Inches(12.333), height=Inches(2.8),
        rows=6, cols=4, data=data,
        font_size=16,
        col_widths=[Inches(1.5), Inches(5.5), Inches(3.0), Inches(2.333)]
    )

    # Honest note about schedule
    add_textbox(
        slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.6),
        "Honest Assessment: We are behind our original schedule.",
        font_size=20, color=RED, bold=True
    )

    items = [
        "Original plan had Task 3 complete by Feb 1 -- revised to Feb 28",
        "Critical path: Hull design must finalize before mold construction begins",
        "Mitigation: Increased meeting frequency to 3x/week; parallel-tracking mix optimization",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(5.2), Inches(11.5), Inches(2.0),
        items, font_size=18, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Jon")
    add_slide_number(slide, 3)
    add_notes(slide, "Task 1-2 complete. Task 3 hull design at 60%. We ARE behind our original schedule but we've adjusted the timeline. Critical path is hull design finalization.")


def build_slide_04_research(prs):
    """Slide 4: Task 1 -- Project Research"""
    print("  Building Slide 4: Research")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 1.1: Project Research")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "ASCE 2026 RFP reviewed cover to cover by all team members",
        "Key dimensional constraints: maximum 24 ft length, must accommodate 4 paddlers",
        "Minimum 6 inches of freeboard with 4 paddlers aboard",
        "Density limitations: concrete must float (< 62.4 PCF)",
        "Reinforcement: only ASCE-approved fiber or mesh reinforcement",
        "No metallic reinforcement (rebar, wire mesh) permitted",
        "Sustainability scoring included for first time in 2026",
        "Each team member assigned specific rule sections for deep review",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(5.0),
        items, font_size=20, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Kayleigh")
    add_slide_number(slide, 4)
    add_notes(slide, "Every team member read the full ASCE 2026 RFP. Key constraints: 6-inch freeboard minimum, density must be below 62.4 PCF (water), only approved non-metallic reinforcement.")


def build_slide_05_sponsorships(prs):
    """Slide 5: Sponsorships"""
    print("  Building Slide 5: Sponsorships")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 1.2: Sponsorships")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Created personalized outreach videos for each potential sponsor",
        "Highlighted NAU engineering program and competition visibility",
        "Sponsors receive logo placement on canoe and in technical paper",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(2.5),
        items, font_size=20, color=DARK_GRAY
    )

    # Sponsor highlight boxes
    for i, (name, status) in enumerate([
        ("Oxendale Auto Group", "CONFIRMED"),
        ("CSEM", "CONFIRMED"),
        ("Additional Sponsors", "In Progress"),
    ]):
        left = Inches(1.0 + i * 4.0)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, Inches(4.2), Inches(3.5), Inches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY if status == "CONFIRMED" else LIGHT_GRAY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE if status == "CONFIRMED" else DARK_GRAY
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = status
        p2.font.size = Pt(16)
        p2.font.color.rgb = GOLD if status == "CONFIRMED" else MEDIUM_GRAY
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Kayleigh")
    add_slide_number(slide, 5)
    add_notes(slide, "We created personalized video pitches for sponsors. Oxendale Auto Group and CSEM have confirmed. We're continuing outreach.")


def build_slide_06_safety(prs):
    """Slide 6: Lab Safety"""
    print("  Building Slide 6: Safety")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 1.3: Lab Safety")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Compiled comprehensive safety binder with all SDS sheets",
        "Met with lab manager to review emergency procedures",
        "All team members completed NAU lab safety training",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(2.0),
        items, font_size=20, color=DARK_GRAY
    )

    # Hazards box
    add_textbox(
        slide, Inches(0.7), Inches(3.5), Inches(5.5), Inches(0.5),
        "Identified Hazards:", font_size=20, color=NAVY, bold=True
    )

    hazards = [
        {"text": "Alkaline burns from fresh concrete (pH ~12-13)", "color": RED, "bold": True},
        {"text": "Silica dust inhalation (perlite, fly ash)", "color": RED, "bold": True},
        {"text": "Heavy lifting: mold sections weigh 50+ lbs", "color": RED, "bold": True},
        {"text": "Eye protection required during mixing and testing", "color": DARK_GRAY},
    ]
    add_bullet_list(
        slide, Inches(0.9), Inches(4.0), Inches(5.5), Inches(2.5),
        hazards, font_size=18, color=DARK_GRAY
    )

    # PPE box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE0)
    shape.line.color.rgb = GOLD

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Required PPE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    ppe_items = [
        "Safety glasses / goggles",
        "Nitrile gloves",
        "N95 dust masks (perlite mixing)",
        "Closed-toe shoes",
        "Long sleeves during pours",
    ]
    for item in ppe_items:
        p2 = tf.add_paragraph()
        p2.text = f"\u2022 {item}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = DARK_GRAY
        p2.space_after = Pt(6)

    add_speaker_tag(slide, "Kayleigh")
    add_slide_number(slide, 6)
    add_notes(slide, "Safety is paramount. Main hazards: alkaline burns from wet concrete, silica dust from perlite and fly ash, heavy lifting. All members trained, safety binder on-site.")


def build_slide_07_materials(prs):
    """Slide 7: Material Research (KEY SLIDE)"""
    print("  Building Slide 7: Materials")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 2.1: Material Research", subtitle_text="KEY SLIDE -- Aggregate Selection")
    add_gold_accent_line(slide, Inches(1.25))

    # Explanation text
    add_textbox(
        slide, Inches(0.5), Inches(1.45), Inches(12.0), Inches(0.5),
        "Specific gravity is the key metric: lower SG = lighter canoe = easier to float",
        font_size=18, color=NAVY, bold=True
    )

    # Aggregate table
    data = [
        ["Aggregate", "Specific Gravity", "Type", "Status"],
        ["Perlite", "0.125", "Volcanic glass", "SELECTED"],
        ["K1 Microspheres", "0.125", "Cenospheres", "SELECTED"],
        ["Poraver", "0.296", "Expanded glass", "CANDIDATE"],
        ["Pumice", "0.641", "Volcanic rock", "EVALUATED"],
        ["Expanded Shale", "1.15", "Lightweight", "REJECTED"],
        ["Natural Sand", "2.65", "Standard", "BASELINE"],
        ["Limestone", "2.70", "Standard", "BASELINE"],
        ["Gravel", "2.75", "Standard", "BASELINE"],
    ]
    table_shape = add_table(
        slide,
        left=Inches(0.5), top=Inches(2.1),
        width=Inches(12.333), height=Inches(4.2),
        rows=9, cols=4, data=data,
        font_size=16,
        col_widths=[Inches(3.5), Inches(3.0), Inches(3.0), Inches(2.833)]
    )

    # Highlight the ultralight rows (1-3) with gold tint
    table = table_shape.table
    for r in range(1, 4):
        for c in range(4):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE0)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True

    # Divider line notation
    add_textbox(
        slide, Inches(0.5), Inches(6.5), Inches(12.0), Inches(0.4),
        "Top 3 aggregates (SG < 1.0) are ultralight -- bottom 4 (SG > 1.0) would make the canoe sink.",
        font_size=14, color=MEDIUM_GRAY, italic=True
    )

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 7)
    add_notes(slide, "Eight aggregates evaluated. The key is specific gravity. Low SG = lighter canoe = easier to float. Perlite and K1 microspheres at 0.125 are our primary lightweight aggregates.")


def build_slide_08_astm_overview(prs):
    """Slide 8: ASTM Testing Overview"""
    print("  Building Slide 8: ASTM Overview")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 2.2: ASTM Testing Standards")
    add_gold_accent_line(slide, Inches(1.25))

    data = [
        ["ASTM Standard", "Test Name", "What It Measures", "Why It Matters"],
        ["C39", "Compressive Strength", "Max load before crushing", "Structural capacity"],
        ["C143", "Slump Test", "Workability / consistency", "Thin-shell placement"],
        ["C138", "Unit Weight", "Wet density (PCF)", "Must float (< 62.4)"],
        ["C136", "Sieve Analysis", "Particle size distribution", "Aggregate grading"],
        ["C496", "Split Tensile", "Indirect tensile strength", "Bottom of canoe in tension"],
        ["C685", "Volumetric Batching", "Consistent mix production", "Repeatability"],
    ]
    add_table(
        slide,
        left=Inches(0.5), top=Inches(1.5),
        width=Inches(12.333), height=Inches(4.0),
        rows=7, cols=4, data=data,
        font_size=15,
        col_widths=[Inches(2.5), Inches(3.0), Inches(3.5), Inches(3.333)]
    )

    add_textbox(
        slide, Inches(0.5), Inches(5.8), Inches(12.0), Inches(0.5),
        "All tests performed following ASTM procedures in NAU Concrete Lab.",
        font_size=16, color=MEDIUM_GRAY, italic=True
    )

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 8)
    add_notes(slide, "Six ASTM test methods guide our material characterization. Each test gives us a specific property we need for the design. C138 density is the most critical -- we need to be below 62.4 PCF.")


def build_slide_09_c39(prs):
    """Slide 9: ASTM C39 -- Compressive Strength"""
    print("  Building Slide 9: C39")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ASTM C39: Compressive Strength")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Standard test for compressive strength of cylindrical concrete specimens",
        "Equipment: Humboldt compression testing machine",
        "Specimen: 4\" diameter x 8\" tall cylinders",
        "Load applied at constant rate until failure",
        "Failure mode: typically cone or shear fracture pattern",
        "Results determine if our mix can handle structural loads",
        "Minimum target: 250 psi at 7 days for thin-shell application",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.5),
        items, font_size=20, color=DARK_GRAY
    )

    # Highlight box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(5.0), Inches(4.3), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Equation"
    p.font.size = Pt(16)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "f'c = P / A"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "(load / cross-section area)"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    p3.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Amit")
    add_slide_number(slide, 9)
    add_notes(slide, "ASTM C39 is the gold standard for compressive strength. We use 4x8 cylinders in the Humboldt machine. Load until failure and compute f'c = P/A.")


def build_slide_10_c143(prs):
    """Slide 10: Slump Test -- ASTM C143"""
    print("  Building Slide 10: C143")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ASTM C143: Slump Test")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Measures workability and consistency of fresh concrete",
        "Fill cone in 3 layers, rod each layer 25 times, lift cone",
        "Measure how much the concrete 'slumps' downward",
        "Target range for canoe: 4-7 inches (fluid enough for thin shell)",
        "Too low (<3\"): won't fill thin mold sections properly",
        "Too high (>8\"): mix may segregate, aggregates settle out",
        "Thin-shell application demands precise workability control",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.5),
        items, font_size=20, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Kayleigh")
    add_slide_number(slide, 10)
    add_notes(slide, "Slump test measures workability. For a thin-shell canoe, we need 4-7 inches of slump. Too stiff and we can't fill the mold; too fluid and aggregates segregate.")


def build_slide_11_c138(prs):
    """Slide 11: Density -- ASTM C138 (KEY SLIDE)"""
    print("  Building Slide 11: C138 Density")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ASTM C138: Unit Weight (Density)", subtitle_text="CRITICAL MEASUREMENT")
    add_gold_accent_line(slide, Inches(1.25))

    # Main point - large
    add_textbox(
        slide, Inches(0.5), Inches(1.5), Inches(12.0), Inches(0.8),
        "THE MOST CRITICAL TEST: Density must be below 62.4 PCF to float.",
        font_size=24, color=RED, bold=True, alignment=PP_ALIGN.CENTER
    )

    items = [
        "Measures wet (fresh) density of concrete in pounds per cubic foot (PCF)",
        "Fill known-volume container, weigh, compute density = weight / volume",
        "Water density = 62.4 PCF -- our concrete MUST be lighter",
        "Target range: 55-65 PCF for lightweight canoe concrete",
        "Lighter is better, but too light sacrifices strength",
        "Every aggregate substitution shifts density -- trade-off with strength",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(2.5), Inches(7.0), Inches(3.5),
        items, font_size=18, color=DARK_GRAY
    )

    # Big threshold box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(2.5), Inches(4.3), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "THRESHOLD"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "62.4 PCF"
    p2.font.size = Pt(44)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Density of Water"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = ""
    p4.font.size = Pt(8)
    p5 = tf.add_paragraph()
    p5.text = "ABOVE = SINKS"
    p5.font.size = Pt(20)
    p5.font.color.rgb = RED
    p5.font.bold = True
    p5.alignment = PP_ALIGN.CENTER
    p6 = tf.add_paragraph()
    p6.text = "BELOW = FLOATS"
    p6.font.size = Pt(20)
    p6.font.color.rgb = RGBColor(0x00, 0xCC, 0x00)
    p6.font.bold = True
    p6.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Amit")
    add_slide_number(slide, 11)
    add_notes(slide, "Density is our most critical measurement. 62.4 PCF is water's density -- we must be below this. Our target is 55-65 PCF: light enough to float, dense enough for strength.")


def build_slide_12_c136(prs):
    """Slide 12: Sieve Analysis -- ASTM C136"""
    print("  Building Slide 12: C136")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ASTM C136: Sieve Analysis")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Determines particle size distribution of aggregates",
        "Stack of progressively finer sieves shaken mechanically",
        "Weigh material retained on each sieve",
        "Well-graded blend: reduces void space between particles",
        "Fewer voids = less cement paste needed = lower density",
        "Critical for our ultralight aggregates (perlite, K1, Poraver)",
        "Sieve sizes: #4, #8, #16, #30, #50, #100, #200, pan",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.5),
        items, font_size=20, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Kayleigh")
    add_slide_number(slide, 12)
    add_notes(slide, "Sieve analysis gives us particle size distribution. A well-graded blend reduces voids, meaning we need less cement paste. This helps keep density low while maintaining packing efficiency.")


def build_slide_13_c496(prs):
    """Slide 13: Tensile Strength -- ASTM C496"""
    print("  Building Slide 13: C496")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ASTM C496: Split Tensile Strength")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Splitting cylinder test: load cylinder on its side",
        "Compressive line load induces transverse tension",
        "Concrete is WEAK in tension -- typically ~10% of compressive strength",
        "The bottom of the canoe hull is in tension under loading",
        "This drives our reinforcement strategy (fiber mesh)",
        "Formula: fct = 2P / (pi * L * D)",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(7.5), Inches(4.0),
        items, font_size=20, color=DARK_GRAY
    )

    # Importance callout
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(2.0), Inches(4.3), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    shape.line.color.rgb = RED
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Why This Matters:"
    p.font.size = Pt(18)
    p.font.color.rgb = RED
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Canoe bottom is in TENSION under paddler weight. Concrete cracks in tension first."
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Amit")
    add_slide_number(slide, 13)
    add_notes(slide, "Concrete is weak in tension, about 10% of its compressive strength. The bottom of the canoe hull is in tension when paddlers sit on it, so tensile strength and fiber reinforcement are critical.")


def build_slide_14_c685(prs):
    """Slide 14: Batching -- ASTM C685"""
    print("  Building Slide 14: C685")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "ASTM C685: Volumetric Batching", subtitle_text="Proof of Work")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        "Volumetric batching ensures consistent, repeatable mixes",
        "Each ingredient measured by volume rather than weight for field practicality",
        "Lab batching procedure documented step-by-step",
        "Dust masks required during perlite and fly ash handling",
        "Mix quantities scaled from lab trial (0.1 CY) to production (0.5 CY)",
        "Graduate Instructor guidance: \"Show proof, not just results\"",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.0),
        items, font_size=20, color=DARK_GRAY
    )

    # Photo placeholder boxes
    for i in range(3):
        left = Inches(1.0 + i * 4.0)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, Inches(5.2), Inches(3.2), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = MEDIUM_GRAY
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        labels = ["[Lab Mixing Photo]", "[Cylinder Casting Photo]", "[Testing Photo]"]
        p.text = labels[i]
        p.font.size = Pt(14)
        p.font.color.rgb = MEDIUM_GRAY
        p.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Kayleigh")
    add_slide_number(slide, 14)
    add_notes(slide, "GI said show proof not just results. This slide shows our actual lab work -- batching, casting cylinders, and testing. Dust masks are essential for perlite handling.")


def build_slide_15_mix_selection(prs):
    """Slide 15: Mix Selection Process"""
    print("  Building Slide 15: Mix Selection")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 2.3: Mix Selection Process")
    add_gold_accent_line(slide, Inches(1.25))

    # Process flow - three boxes with arrows
    box_labels = [
        ("3 Precursor\nMixes", "Literature review\n& material screening"),
        ("3 Trial\nMixes", "Lab-batched\n& ASTM-tested"),
        ("Decision\nMatrix", "Weighted criteria\nscoring"),
    ]
    for i, (title, desc) in enumerate(box_labels):
        left = Inches(0.8 + i * 4.2)
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, Inches(2.0), Inches(3.5), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = ""
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(16)
        p3.font.color.rgb = GOLD
        p3.alignment = PP_ALIGN.CENTER

        # Arrow (except after last)
        if i < 2:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(4.5 + i * 4.2), Inches(2.8), Inches(0.7), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()

    # Summary text below
    items = [
        "Systematic approach: research-based precursors -> lab testing -> data-driven selection",
        "Each mix tested per ASTM C39, C143, C138, C496",
        "Decision matrix weighted by project priorities (density 30%, strength 25%, workability 20%, cost 15%, sustainability 10%)",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(5.0), Inches(11.5), Inches(2.0),
        items, font_size=18, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 15)
    add_notes(slide, "Our process: three precursor mixes from literature, three lab trial mixes, then a weighted decision matrix to select the best. This ensures a data-driven, defensible choice.")


def build_slide_16_mix1(prs):
    """Slide 16: Mix 1 Results"""
    print("  Building Slide 16: Mix 1")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Mix 1: High-Strength Baseline")
    add_gold_accent_line(slide, Inches(1.25))

    # Results table
    data = [
        ["Property", "Result", "Target", "Verdict"],
        ["Compressive Strength (7-day)", "3,116 psi", "> 250 psi", "PASS"],
        ["Tensile Strength", "385 psi", "> 40 psi", "PASS"],
        ["Density", "78.4 PCF", "< 62.4 PCF", "FAIL"],
        ["Slump", "1.5\"", "4-7\"", "MARGINAL"],
    ]
    table_shape = add_table(
        slide,
        left=Inches(0.5), top=Inches(1.5),
        width=Inches(12.333), height=Inches(2.5),
        rows=5, cols=4, data=data,
        font_size=18,
        col_widths=[Inches(4.0), Inches(3.0), Inches(3.0), Inches(2.333)]
    )

    # Highlight the FAIL cell
    table = table_shape.table
    fail_cell = table.cell(3, 3)
    fail_cell.fill.solid()
    fail_cell.fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    for p in fail_cell.text_frame.paragraphs:
        p.font.color.rgb = RED
        p.font.bold = True

    # Big fail callout
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(4.5), Inches(8.333), Inches(2.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0xE0)
    shape.line.color.rgb = RED
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DENSITY FAILURE: 78.4 PCF > 62.4 PCF"
    p.font.size = Pt(28)
    p.font.color.rgb = RED
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "This mix would SINK. Strong but way too heavy."
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 16)
    add_notes(slide, "Mix 1 was our high-strength baseline. Excellent compressive strength at 3,116 psi, but density of 78.4 PCF is way above the 62.4 threshold. This canoe would sink.")


def build_slide_17_mix2(prs):
    """Slide 17: Mix 2 Results"""
    print("  Building Slide 17: Mix 2")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Mix 2: Slag Cement Addition")
    add_gold_accent_line(slide, Inches(1.25))

    data = [
        ["Property", "Mix 1", "Mix 2", "Change"],
        ["Compressive (7-day)", "3,116 psi", "~2,800 psi", "Slight decrease"],
        ["Density", "78.4 PCF", "83.6 PCF", "WORSE (+5.2)"],
        ["Slump", "1.5\"", "2.0\"", "Improved"],
        ["Additive", "--", "Slag cement", "Added"],
    ]
    table_shape = add_table(
        slide,
        left=Inches(0.5), top=Inches(1.5),
        width=Inches(12.333), height=Inches(2.8),
        rows=5, cols=4, data=data,
        font_size=18,
        col_widths=[Inches(3.5), Inches(3.0), Inches(3.0), Inches(2.833)]
    )

    # Highlight density worse
    table = table_shape.table
    cell = table.cell(2, 3)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = RED
        p.font.bold = True

    items = [
        {"text": "Slag cement improved workability (slump: 1.5\" -> 2.0\")", "color": DARK_GRAY},
        {"text": "But density got WORSE: 83.6 PCF (even heavier than Mix 1)", "color": RED, "bold": True},
        {"text": "Trade-off lesson: slag helps workability but hurts density", "color": DARK_GRAY},
        {"text": "Key insight: must change AGGREGATES, not just cementitious materials", "color": NAVY, "bold": True},
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(4.7), Inches(11.5), Inches(2.5),
        items, font_size=20, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 17)
    add_notes(slide, "Mix 2 added slag cement. Workability improved slightly, but density got WORSE at 83.6 PCF. This taught us that to fix density, we need to change the aggregates, not just the cement.")


def build_slide_18_mix3(prs):
    """Slide 18: Mix 3 Results (MOST IMPORTANT)"""
    print("  Building Slide 18: Mix 3 (Breakthrough)")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Mix 3: Lightweight Breakthrough", subtitle_text="MOST IMPORTANT RESULT")
    add_gold_accent_line(slide, Inches(1.25))

    # Key changes
    add_textbox(
        slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
        "Key Changes from Mix 1/2:", font_size=20, color=NAVY, bold=True
    )
    changes = [
        "Portland cement cut in HALF: 97 lbs vs 209 lbs",
        "Added fly ash as supplementary cementitious material",
        "Increased perlite and K1 microsphere content",
        "Result: dramatically lower density",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(2.0),
        changes, font_size=18, color=DARK_GRAY
    )

    # Results table
    data = [
        ["Property", "Result", "Target"],
        ["Density", "58.6 PCF", "< 62.4 PCF"],
        ["Compressive (7-day)", "307 psi", "> 250 psi"],
        ["Slump", "6.0\"", "4-7\""],
        ["Cost", "$1,200/CY", "Minimize"],
    ]
    table_shape = add_table(
        slide,
        left=Inches(0.5), top=Inches(4.0),
        width=Inches(6.0), height=Inches(2.5),
        rows=5, cols=3, data=data,
        font_size=16,
        col_widths=[Inches(2.5), Inches(2.0), Inches(1.5)]
    )

    # Highlight density PASS row
    table = table_shape.table
    density_cell = table.cell(1, 1)
    density_cell.fill.solid()
    density_cell.fill.fore_color.rgb = RGBColor(0xD0, 0xFF, 0xD0)
    for p in density_cell.text_frame.paragraphs:
        p.font.color.rgb = GREEN
        p.font.bold = True

    # Big success callout on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xFF, 0xE0)
    shape.line.color.rgb = GREEN
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "IT FLOATS!"
    p.font.size = Pt(36)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = ""
    p3 = tf.add_paragraph()
    p3.text = "58.6 PCF < 62.4 PCF"
    p3.font.size = Pt(28)
    p3.font.color.rgb = NAVY
    p3.font.bold = True
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = ""
    p5 = tf.add_paragraph()
    p5.text = "Density below water threshold"
    p5.font.size = Pt(18)
    p5.font.color.rgb = DARK_GRAY
    p5.alignment = PP_ALIGN.CENTER

    # Mix composition image
    img_path = os.path.join(FIGURES_DIR, "mix_composition_comparison.png")
    add_image_safe(slide, img_path, Inches(7.0), Inches(4.6), width=Inches(5.8), height=Inches(2.5))

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 18)
    add_notes(slide, "Our breakthrough. Portland cement cut in half, added fly ash, increased perlite and K1. Density dropped to 58.6 PCF -- below 62.4. This mix FLOATS. Compressive at 307 psi meets our minimum. 6-inch slump is excellent workability.")


def build_slide_19_decision_matrix(prs):
    """Slide 19: Decision Matrix"""
    print("  Building Slide 19: Decision Matrix")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 2.4: Decision Matrix")
    add_gold_accent_line(slide, Inches(1.25))

    items = [
        {"text": "Density weighted 30% -- floating is non-negotiable", "color": NAVY, "bold": True},
        "Compressive strength: 25%",
        "Workability: 20%",
        "Cost: 15%",
        "Sustainability: 10%",
    ]
    add_bullet_list(
        slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.5),
        items, font_size=18, color=DARK_GRAY
    )

    # Scores summary
    data = [
        ["Mix", "Score", "Rank"],
        ["Mix 1: Baseline", "65.20", "3rd"],
        ["Mix 2: Slag", "73.53", "1st"],
        ["Mix 3: Lightweight", "70.07", "2nd"],
    ]
    add_table(
        slide,
        left=Inches(0.5), top=Inches(4.2),
        width=Inches(5.5), height=Inches(2.0),
        rows=4, cols=3, data=data,
        font_size=18,
        col_widths=[Inches(2.5), Inches(1.5), Inches(1.5)]
    )

    # Decision matrix image on right
    img_path = os.path.join(FIGURES_DIR, "mix_decision_matrix.png")
    add_image_safe(slide, img_path, Inches(6.5), Inches(1.5), width=Inches(6.3), height=Inches(5.0))

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 19)
    add_notes(slide, "Decision matrix with density weighted at 30% because floating is non-negotiable. Mix 2 scores highest at 73.53, Mix 3 close at 70.07. Both have strengths we want to combine.")


def build_slide_20_selected_mix(prs):
    """Slide 20: Mix Selection Decision"""
    print("  Building Slide 20: Selected Mix")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Selected Approach: Optimized Mix 2 + Mix 3 Aggregates")
    add_gold_accent_line(slide, Inches(1.25))

    # Strategy box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.6), Inches(12.333), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE0)
    shape.line.color.rgb = GOLD
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategy: Combine Mix 2's cementitious system with Mix 3's aggregate ratios"
    p.font.size = Pt(22)
    p.font.color.rgb = NAVY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = ""
    p3 = tf.add_paragraph()
    p3.text = "Mix 2 brings workability and strength from slag cement"
    p3.font.size = Pt(18)
    p3.font.color.rgb = DARK_GRAY
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = "Mix 3 brings the lightweight aggregate ratios that achieved 58.6 PCF"
    p4.font.size = Pt(18)
    p4.font.color.rgb = DARK_GRAY
    p4.alignment = PP_ALIGN.CENTER

    # Going forward items
    add_textbox(
        slide, Inches(0.5), Inches(4.5), Inches(12.0), Inches(0.5),
        "Next Steps for Mix Optimization:", font_size=22, color=NAVY, bold=True
    )

    items = [
        "Target: bring density below 62.4 PCF while maintaining strength > 250 psi",
        "Optimize perlite/K1 ratio for best density-strength trade-off",
        "Test 28-day strength (expect significant gains over 7-day)",
        "Scale from 0.1 CY trial batch to 0.5 CY production batch",
        "Finalize mix design by March 1 for construction start",
    ]
    add_bullet_list(
        slide, Inches(0.7), Inches(5.0), Inches(11.5), Inches(2.0),
        items, font_size=18, color=DARK_GRAY
    )

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 20)
    add_notes(slide, "Going forward: optimize Mix 2 by incorporating Mix 3's aggregate ratios. The goal is to get the best of both -- workability from slag cement and low density from ultralight aggregates.")


def build_slide_21_hull_design(prs):
    """Slide 21: Hull Design -- Design Criteria (CRITICAL SLIDE)"""
    print("  Building Slide 21: Hull Design (Critical)")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Task 3.1: Hull Design -- Design Criteria & Preliminary Results",
                  subtitle_text="CRITICAL SLIDE")
    add_gold_accent_line(slide, Inches(1.25))

    # LEFT COLUMN: ASCE Requirements
    add_textbox(
        slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
        "ASCE 2026 Requirements", font_size=20, color=NAVY, bold=True
    )

    req_data = [
        ["Parameter", "Requirement"],
        ["Freeboard (4-person)", ">= 6.0\""],
        ["Metacentric Height (GM)", ">= 6.0\""],
        ["Safety Factor", ">= 2.0 (ACI 318-25 LRFD)"],
        ["Punching Shear DCR", "< 1.0 (ACI 22.6.5.2)"],
        ["Density", "< 62.4 PCF"],
    ]
    add_table(
        slide,
        left=Inches(0.3), top=Inches(2.0),
        width=Inches(5.8), height=Inches(2.8),
        rows=6, cols=2, data=req_data,
        font_size=15,
        col_widths=[Inches(2.8), Inches(3.0)]
    )

    # RIGHT COLUMN: Results
    add_textbox(
        slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
        "Candidate Design Results", font_size=20, color=NAVY, bold=True
    )

    res_data = [
        ["Parameter", "Value", "Status"],
        ["Dimensions", "216\" x 36\" x 18\"", "--"],
        ["Thickness", "0.75\"", "--"],
        ["Weight", "~318 lbs", "--"],
        ["Freeboard", "12.82\"", "PASS"],
        ["Metacentric Ht", "15.34\"", "PASS"],
        ["Safety Factor", "6.39", "PASS"],
        ["Punch DCR", "0.189", "PASS"],
    ]
    table_shape = add_table(
        slide,
        left=Inches(6.8), top=Inches(2.0),
        width=Inches(6.0), height=Inches(3.5),
        rows=8, cols=3, data=res_data,
        font_size=15,
        col_widths=[Inches(2.5), Inches(2.0), Inches(1.5)]
    )

    # Green PASS cells
    table = table_shape.table
    for r in [4, 5, 6, 7]:
        cell = table.cell(r, 2)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xD0, 0xFF, 0xD0)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = GREEN
            p.font.bold = True

    # SolidWorks image below - spanning both columns
    img_path = os.path.join(FIGURES_DIR, "solidworks_3d_dimensioned.png")
    add_image_safe(slide, img_path, Inches(0.3), Inches(5.6), width=Inches(5.5), height=Inches(1.6))

    # FBD image
    img_path2 = os.path.join(FIGURES_DIR, "report_fig1_fbd_buoyancy.png")
    add_image_safe(slide, img_path2, Inches(6.3), Inches(5.6), width=Inches(4.0), height=Inches(1.6))

    add_footer_text(slide, "Structural analysis per ACI 318-25 LRFD. Cross-checked: Python v2.1 + Excel")

    add_speaker_tag(slide, "Trevion")
    add_slide_number(slide, 21)
    add_notes(slide, "ASCE requires 6-inch minimum freeboard and adequate stability. Our 216x36x18 design gives 12.82 inches of freeboard -- double the minimum -- and a safety factor of 6.39. Full shear/moment diagrams at 60% review.")


def build_slide_22_structural_preview(prs):
    """Slide 22: Structural Analysis Preview"""
    print("  Building Slide 22: Structural Preview")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Structural Analysis -- FBD & Section Properties")
    add_gold_accent_line(slide, Inches(1.25))

    # Main FBD image
    img_path = os.path.join(FIGURES_DIR, "report_fig1_fbd_buoyancy.png")
    add_image_safe(slide, img_path, Inches(0.5), Inches(1.5), width=Inches(8.0), height=Inches(4.5))

    # MOI image on right
    img_path2 = os.path.join(FIGURES_DIR, "report_fig2_moi_rectangles.png")
    add_image_safe(slide, img_path2, Inches(8.8), Inches(1.5), width=Inches(4.0), height=Inches(3.5))

    # Text below MOI
    add_textbox(
        slide, Inches(8.8), Inches(5.2), Inches(4.0), Inches(0.8),
        "Parallel axis theorem applied\nto U-shell cross section",
        font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, italic=True
    )

    # Footer note
    add_textbox(
        slide, Inches(0.5), Inches(6.3), Inches(8.0), Inches(0.5),
        "Full structural analysis with shear/moment diagrams at 60% review",
        font_size=18, color=NAVY, bold=True
    )

    add_speaker_tag(slide, "Trevion")
    add_slide_number(slide, 22)
    add_notes(slide, "Free body diagram showing distributed loading, buoyancy, and crew weight. Section properties computed using parallel axis theorem for the U-shell cross section. Full analysis coming at 60% review.")


def build_slide_23_3d_hull(prs):
    """Slide 23: 3D Hull Model"""
    print("  Building Slide 23: 3D Hull")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "3D Hull Model -- SolidWorks Dimensions")
    add_gold_accent_line(slide, Inches(1.25))

    # SolidWorks image - large, centered
    img_path = os.path.join(FIGURES_DIR, "solidworks_3d_dimensioned.png")
    add_image_safe(slide, img_path, Inches(0.5), Inches(1.5), width=Inches(12.333), height=Inches(5.5))

    add_speaker_tag(slide, "Trevion")
    add_slide_number(slide, 23)
    add_notes(slide, "SolidWorks 3D model showing full dimensions of our 216x36x18 inch hull design. The model is used for volume calculations, mold design, and construction drawings.")


def build_slide_24_references(prs):
    """Slide 24: References"""
    print("  Building Slide 24: References")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "References")
    add_gold_accent_line(slide, Inches(1.25))

    references = [
        "[1]  ACI 318-25, Building Code Requirements for Structural Concrete, ACI, 2025",
        "[2]  ASCE 2026 Concrete Canoe Competition Rules & Regulations",
        "[3]  ASTM C39/C39M-21, Standard Test Method for Compressive Strength",
        "[4]  ASTM C143/C143M-20, Standard Test Method for Slump",
        "[5]  ASTM C138/C138M-17a, Standard Test Method for Unit Weight",
        "[6]  ASTM C136/C136M-19, Standard Test Method for Sieve Analysis",
        "[7]  ASTM C496/C496M-17, Standard Test Method for Splitting Tensile",
        "[8]  ASTM C685/C685M-17, Standard Specification for Volumetric Batching",
        "[9]  Lewis, E.V., Principles of Naval Architecture, SNAME, 1988",
        "[10] Beer, Johnston, DeWolf, Mechanics of Materials, 8th Ed.",
    ]

    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(references):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(8)

    add_speaker_tag(slide, "Theo")
    add_slide_number(slide, 24)
    add_notes(slide, "References for all standards, codes, and textbooks cited in this presentation. ACI 318-25 is the primary structural code. ASCE 2026 RFP governs all competition requirements.")


# =============================================================================
# MAIN
# =============================================================================
def main():
    print("=" * 60)
    print("Generating Pluto Jacks 30% Capstone Review Presentation")
    print("=" * 60)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all 24 slides
    build_slide_01_title(prs)
    build_slide_02_introduction(prs)
    build_slide_03_schedule(prs)
    build_slide_04_research(prs)
    build_slide_05_sponsorships(prs)
    build_slide_06_safety(prs)
    build_slide_07_materials(prs)
    build_slide_08_astm_overview(prs)
    build_slide_09_c39(prs)
    build_slide_10_c143(prs)
    build_slide_11_c138(prs)
    build_slide_12_c136(prs)
    build_slide_13_c496(prs)
    build_slide_14_c685(prs)
    build_slide_15_mix_selection(prs)
    build_slide_16_mix1(prs)
    build_slide_17_mix2(prs)
    build_slide_18_mix3(prs)
    build_slide_19_decision_matrix(prs)
    build_slide_20_selected_mix(prs)
    build_slide_21_hull_design(prs)
    build_slide_22_structural_preview(prs)
    build_slide_23_3d_hull(prs)
    build_slide_24_references(prs)

    # Save
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"\nPresentation saved to: {OUTPUT_FILE}")
    print(f"Total slides: {len(prs.slides)}")
    print("Done!")


if __name__ == "__main__":
    main()
