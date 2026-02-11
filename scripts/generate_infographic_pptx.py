#!/usr/bin/env python3
"""
NAU ASCE Concrete Canoe 2026 — Editable PowerPoint Infographic
Single 24"x18" slide. Every element is click-drag-resize editable.
Modern visual design: big metric circles, embedded figures, photo placeholders.
~168 words total — visuals dominate.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# CONSTANTS
# ============================================================
SLIDE_W = Inches(24)
SLIDE_H = Inches(18)

# NAU Brand Colors
NAVY       = RGBColor(0x00, 0x34, 0x66)
NAVY_DARK  = RGBColor(0x00, 0x22, 0x44)
GOLD       = RGBColor(0xFF, 0xB8, 0x1C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF4, 0xF6, 0xF9)
GREEN      = RGBColor(0x1B, 0x7D, 0x3A)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)

FONT = "Calibri"
FIG_DIR = "/root/concrete-canoe-project2026/reports/figures"
OUT_DIR = "/root/concrete-canoe-project2026/reports"


# ============================================================
# HELPERS
# ============================================================
def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=DARK_GRAY, bold=False,
                alignment=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = alignment
    return tb


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=None, line_color=None, line_width=None):
    s = slide.shapes.add_shape(shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        if line_width:
            s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def add_image_safe(slide, filename, left, top, width=None, height=None):
    path = os.path.join(FIG_DIR, filename)
    if not os.path.exists(path):
        # Draw a placeholder rectangle if image missing
        s = add_shape(slide, MSO_SHAPE.RECTANGLE, left, top,
                      width or 3, height or 2, fill_color=LIGHT_BG, line_color=NAVY, line_width=1)
        s.text_frame.paragraphs[0].text = f"[{filename}]"
        s.text_frame.paragraphs[0].font.size = Pt(10)
        s.text_frame.paragraphs[0].font.color.rgb = MED_GRAY
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return s
    kwargs = {}
    if width: kwargs['width'] = Inches(width)
    if height: kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


def add_metric_circle(slide, left, top, size, big_number, label):
    """Navy oval with gold border, big white number, small gold label."""
    s = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size,
                  fill_color=NAVY, line_color=GOLD, line_width=3)

    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = big_number
    p1.font.size = Pt(int(size * 12))  # scale font with circle
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = FONT
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(0)

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(int(size * 4.5))
    p2.font.color.rgb = GOLD
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)

    return s


def add_photo_placeholder(slide, left, top, width, height, label, sublabel):
    """Dashed-border box for user to drop in photos."""
    s = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                  fill_color=RGBColor(0xF5, 0xF5, 0xF5), line_color=NAVY, line_width=1.5)

    # Rounded corners
    if s.adjustments:
        s.adjustments[0] = 0.08

    # Dashed border via XML
    spPr = s._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = s._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')

    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.font.name = FONT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = sublabel
    p2.font.size = Pt(10)
    p2.font.color.rgb = MED_GRAY
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    return s


def add_section_header(slide, left, top, text, line_w=3.0):
    """Bold header text + gold underline bar."""
    add_textbox(slide, left, top, line_w + 2, 0.5, text,
                font_size=20, color=NAVY, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top + 0.4, line_w, 0.06,
              fill_color=GOLD)


def add_compliance_row(slide, left, top, requirement, value):
    """Green checkmark circle + requirement + value."""
    # Green circle with checkmark
    s = add_shape(slide, MSO_SHAPE.OVAL, left, top, 0.45, 0.45,
                  fill_color=GREEN)
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "\u2713"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    # Requirement text
    add_textbox(slide, left + 0.6, top + 0.02, 4.5, 0.4, requirement,
                font_size=13, color=DARK_GRAY)

    # Value text (bold)
    add_textbox(slide, left + 5.2, top + 0.02, 2.0, 0.4, value,
                font_size=13, color=NAVY, bold=True)


def add_kv_row(slide, left, top, label, value, label_w=2.5, font_size=13, bold_val=False):
    """Key-value row for specs."""
    add_textbox(slide, left, top, label_w, 0.35, label,
                font_size=font_size, color=MED_GRAY)
    add_textbox(slide, left + label_w, top, 3.5, 0.35, value,
                font_size=font_size, color=NAVY, bold=bold_val)


# ============================================================
# BUILD SECTIONS
# ============================================================
def build_header(slide):
    # Navy gradient background
    hdr = add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 24, 2.0,
                    fill_color=NAVY_DARK)

    # Gold accent line at bottom of header
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 1.95, 24, 0.06, fill_color=GOLD)

    # Title
    add_textbox(slide, 0.5, 0.2, 23, 1.0, "PLUTO JACKS",
                font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, 0.5, 1.15, 23, 0.6,
                "NAU ASCE Concrete Canoe 2026  \u2014  Design A  |  Northern Arizona University",
                font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)


def build_metric_circles(slide):
    metrics = [
        (0.6, 2.5,  3.3, "174", "LBS TOTAL"),
        (4.2, 2.5,  3.3, "2.30", "SAFETY FACTOR"),
        (0.6, 6.2,  3.3, '11.4"', "FREEBOARD"),
        (4.2, 6.2,  3.3, '8.68"', "GM HEIGHT"),
    ]
    for left, top, size, num, label in metrics:
        add_metric_circle(slide, left, top, size, num, label)

    # Two smaller circles for secondary metrics
    add_metric_circle(slide, 0.8, 9.9, 2.6, "1,500", "PSI FLEXURAL")
    add_metric_circle(slide, 3.7, 9.9, 2.6, "60", "PCF DRY DENSITY")


def build_hero_image(slide):
    # Navy frame behind image
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.0, 2.2, 8.8, 7.2,
              fill_color=NAVY)

    # 3D hull render
    add_image_safe(slide, "3d_hull_design_A.png", 8.15, 2.35, width=8.5, height=6.5)

    # Dimension strip overlay at bottom of hero
    strip = add_shape(slide, MSO_SHAPE.RECTANGLE, 8.15, 8.2, 8.5, 0.7,
                      fill_color=NAVY)
    # Make it semi-transparent via XML
    try:
        spPr = strip._element.find(qn('p:spPr'))
        solidFill = spPr.find(qn('a:solidFill'))
        srgbClr = solidFill.find(qn('a:srgbClr'))
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '80000')
    except Exception:
        pass

    add_textbox(slide, 8.3, 8.25, 8.2, 0.6,
                '192" (16.0 ft) L  \u2502  32" Beam  \u2502  17" Depth  \u2502  0.5" Wall',
                font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def build_mix_specs(slide):
    add_section_header(slide, 17.2, 2.3, "MIX DESIGN", line_w=2.5)

    # Background panel
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 17.1, 2.8, 6.5, 3.4,
              fill_color=LIGHT_BG)

    specs = [
        ("Wet Density", "~80 PCF"),
        ("Dry Density", "60 PCF"),
        ("Slump", "4 - 6 in."),
        ("Air Content", "8 - 12%"),
        ("Compressive", "2,000 psi"),
        ("Flexural", "1,500 psi"),
        ("w/cm Ratio", "0.35"),
    ]
    for i, (k, v) in enumerate(specs):
        add_kv_row(slide, 17.4, 3.0 + i * 0.42, k, v, label_w=2.2, font_size=13, bold_val=True)


def build_photo_placeholders(slide):
    add_textbox(slide, 17.2, 6.35, 6.4, 0.4, "YOUR PHOTOS",
                font_size=14, color=NAVY, bold=True)

    placeholders = [
        (17.2, 6.8, "CYLINDER TEST", "Insert split cylinder photo"),
        (20.4, 6.8, "FLOAT TEST", "Insert canoe on water"),
        (17.2, 8.7, "CONSTRUCTION", "Insert layup / mold photo"),
        (20.4, 8.7, "TEAM PHOTO", "Insert team photo here"),
    ]
    for left, top, label, sub in placeholders:
        add_photo_placeholder(slide, left, top, 3.0, 1.7, label, sub)


def build_hull_cross_section(slide):
    add_section_header(slide, 0.4, 12.8, "HULL CROSS-SECTION", line_w=3.5)
    add_image_safe(slide, "cross_section_design_A.png", 0.6, 13.2, width=7.0, height=2.2)


def build_asce_compliance(slide):
    add_section_header(slide, 8.4, 10.6, "ASCE COMPLIANCE", line_w=3.0)

    checks = [
        ('Freeboard  \u2265 6.0"', '11.4"'),
        ('Metacentric Height  \u2265 6.0"', '8.68"'),
        ('Safety Factor  \u2265 2.0', '2.30'),
        ('Canoe Weight  \u2264 237 lb', '174.3 lb'),
        ('Cement Ratio  \u2264 0.40', '0.35'),
        ('Reinforcement POA  \u2265 40%', '42%'),
    ]
    for i, (req, val) in enumerate(checks):
        add_compliance_row(slide, 8.4, 11.2 + i * 0.58, req, val)

    # "6/6 PASS" banner
    banner = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.5, 14.8, 4.5, 0.65,
                       fill_color=GREEN)
    tf = banner.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "6 / 6  CHECKS PASS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


def build_weight_breakdown(slide):
    add_section_header(slide, 16.4, 10.6, "WEIGHT BREAKDOWN", line_w=3.5)

    # Big number
    add_textbox(slide, 16.5, 11.2, 3.0, 1.5, "174.3",
                font_size=52, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 16.5, 12.6, 3.0, 0.5, "lbs total",
                font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Horizontal bars (proportional)
    max_w = 4.2
    bars = [
        ("Concrete Shell", 163.1, NAVY),
        ("PVA Mesh", 8.2, GOLD),
        ("Finish/Sealant", 3.0, RGBColor(0x8A, 0xAD, 0xCF)),
    ]
    bar_x = 19.5
    for i, (name, weight, color) in enumerate(bars):
        bar_w = max(max_w * (weight / 174.3), 0.3)  # min visible width
        bar_y = 11.3 + i * 0.85

        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, bar_w, 0.55,
                  fill_color=color)

        pct = weight / 174.3 * 100
        add_textbox(slide, bar_x, bar_y + 0.58, bar_w + 1.5, 0.35,
                    f"{name}  {weight} lbs ({pct:.1f}%)",
                    font_size=10, color=DARK_GRAY)


def build_construction_strip(slide):
    # Thin divider line
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, 15.4, 23.2, 0.04, fill_color=GOLD)

    add_textbox(slide, 0.4, 15.45, 6.0, 0.4, "CONSTRUCTION SEQUENCE",
                font_size=16, color=NAVY, bold=True)

    steps = [
        ("construction_step1_mold_build.png", "1. Mold"),
        ("construction_step2_mold_surface.png", "2. Surface"),
        ("construction_step3_reinforcement.png", "3. Reinf."),
        ("construction_step4_concrete_application.png", "4. Concrete"),
        ("construction_step5_curing.png", "5. Curing"),
        ("construction_step6_demold.png", "6. Demold"),
        ("construction_step7_finishing.png", "7. Finish"),
    ]
    for i, (filename, label) in enumerate(steps):
        x = 0.5 + i * 3.3
        add_image_safe(slide, filename, x, 15.85, width=3.0, height=1.05)
        add_textbox(slide, x, 16.92, 3.0, 0.25, label,
                    font_size=10, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)


def build_innovation_badges(slide):
    """Small gold accent circles as visual tags."""
    badges = [
        (7.8, 4.5, "PORAVER\nGLASS"),
        (7.8, 6.0, "CO\u2082\nCURED"),
        (7.8, 7.5, "AI\nOPTIMIZED"),
    ]
    for left, top, text in badges:
        s = add_shape(slide, MSO_SHAPE.OVAL, left, top, 1.1, 1.1,
                      fill_color=GOLD)
        tf = s.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = NAVY_DARK
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER


def build_footer(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 17.2, 24, 0.8, fill_color=NAVY_DARK)

    add_textbox(slide, 0.5, 17.25, 23, 0.35,
                "Northern Arizona University  |  ASCE Student Chapter  |  Flagstaff, Arizona",
                font_size=13, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 17.55, 23, 0.3,
                "2026 National Concrete Canoe Competition",
                font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# MAIN
# ============================================================
def build_infographic():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Build all sections (z-order: first = back, last = front)
    print("  Building header...")
    build_header(slide)
    print("  Building hero image...")
    build_hero_image(slide)
    print("  Building metric circles...")
    build_metric_circles(slide)
    print("  Building mix specs...")
    build_mix_specs(slide)
    print("  Building photo placeholders...")
    build_photo_placeholders(slide)
    print("  Building hull cross-section...")
    build_hull_cross_section(slide)
    print("  Building ASCE compliance...")
    build_asce_compliance(slide)
    print("  Building weight breakdown...")
    build_weight_breakdown(slide)
    print("  Building construction strip...")
    build_construction_strip(slide)
    print("  Building innovation badges...")
    build_innovation_badges(slide)
    print("  Building footer...")
    build_footer(slide)

    # Save
    out_path = os.path.join(OUT_DIR, "infographic_design_A.pptx")
    prs.save(out_path)
    size_kb = os.path.getsize(out_path) / 1024
    print(f"\nSaved: {out_path} ({size_kb:.0f} KB)")
    print("Open in PowerPoint or Google Slides — every element is editable!")


if __name__ == "__main__":
    print("Generating editable PowerPoint infographic...")
    build_infographic()
    print("Done!")
